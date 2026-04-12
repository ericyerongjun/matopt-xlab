import base64
import io
import json
import os
import re
import shutil
import subprocess
import textwrap
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any, Dict, List, Literal, Optional
from xml.etree import ElementTree as ET

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from openai import APIConnectionError, APIStatusError, APITimeoutError, OpenAI
from pydantic import BaseModel, Field

try:
    import fitz  # type: ignore
except Exception:  # pragma: no cover
    fitz = None

try:
    from docx import Document  # type: ignore
except Exception:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None

from PIL import Image, ImageDraw, ImageFont

from agentic import run_agentic_chat
from skill_router import resolve_skill_selection
from skill_store import delete_skill, list_skills, read_skill, slugify, upsert_skill


class ChatMessage(BaseModel):
    role: Literal["system", "user", "assistant"]
    content: str


class VisualArtifact(BaseModel):
    type: Literal["svg", "interactive_html", "plotly"]
    title: str
    description: str = ""
    code: str


class ChatRequest(BaseModel):
    messages: List[ChatMessage]
    model: Optional[str] = None
    stream: Optional[bool] = False
    agentic: Optional[bool] = True
    skill_slugs: List[str] = Field(default_factory=list)


class ChatResponse(BaseModel):
    id: str
    content: str
    tool_calls: List[dict] = Field(default_factory=list)
    artifacts: List[VisualArtifact] = Field(default_factory=list)
    metadata: Dict[str, Any] = Field(default_factory=dict)
    usage: Optional[dict] = None


class DocumentParseResponse(BaseModel):
    markdown: str
    latex_blocks: List[str] = Field(default_factory=list)
    metadata: Dict[str, Any] = Field(default_factory=dict)
    preview_image_data_url: Optional[str] = None


class SkillUpsertRequest(BaseModel):
    slug: str
    title: Optional[str] = None
    content: str = ""


class SkillResponse(BaseModel):
    slug: str
    title: str
    content: str
    path: str
    updated_at: str
    size_bytes: int


class WorkflowDownload(BaseModel):
    url: str
    filename: str
    size_bytes: int


class PdfWorkflowResponse(BaseModel):
    assistant_markdown: str
    analysis: str
    selected_skills: List[str] = Field(default_factory=list)
    download: WorkflowDownload
    metadata: Dict[str, Any] = Field(default_factory=dict)


app = FastAPI(title="MatOpt Chat Backend", version="0.1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

GENERATED_ROOT = Path(__file__).resolve().parent / "generated"
GENERATED_ROOT.mkdir(parents=True, exist_ok=True)
app.mount("/generated", StaticFiles(directory=str(GENERATED_ROOT)), name="generated")


SUPPORTED_EXTENSIONS = {
    "pdf",
    "pptx",
    "ppt",
    "doc",
    "docx",
    "r",
    "rmd",
    "py",
    "ipynb",
    "c",
    "cpp",
    "java",
    "png",
    "jpg",
    "jpeg",
    "heic",
}

TEXT_EXTENSIONS = {"r", "rmd", "py", "c", "cpp", "java"}
OFFICE_EXTENSIONS = {"doc", "docx", "ppt", "pptx"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "heic"}


@app.exception_handler(Exception)
async def unhandled_exception_handler(_, exc: Exception):
    return JSONResponse(
        status_code=500,
        content={"detail": f"Unhandled backend error: {exc}"},
    )


def _get_openai_api_key() -> str:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY is not set")
    return api_key


def get_client() -> OpenAI:
    api_key = _get_openai_api_key()
    return OpenAI(api_key=api_key)


def _get_extension(filename: str) -> str:
    return Path(filename).suffix.lower().lstrip(".")


def _to_data_url(mime: str, content: bytes) -> str:
    encoded = base64.b64encode(content).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def _mime_for_ext(ext: str) -> str:
    return {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "heic": "image/heic",
    }.get(ext, "application/octet-stream")


def _safe_decode(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_json_object(raw: str) -> Dict[str, Any]:
    text = raw.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    try:
        parsed = json.loads(text)
        return parsed if isinstance(parsed, dict) else {}
    except json.JSONDecodeError:
        pass

    match = re.search(r"\{[\s\S]*\}", text)
    if not match:
        return {}
    try:
        parsed = json.loads(match.group(0))
        return parsed if isinstance(parsed, dict) else {}
    except json.JSONDecodeError:
        return {}


def _extract_script(raw: str) -> str:
    text = raw.strip()
    fenced = re.search(r"```python\s*(.*?)```", text, flags=re.DOTALL | re.IGNORECASE)
    if fenced:
        return fenced.group(1).strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:python)?\s*", "", text, flags=re.IGNORECASE)
        text = re.sub(r"\s*```$", "", text)
    return text.strip()


def _safe_output_filename(stem: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9_-]+", "-", stem).strip("-").lower()
    if not cleaned:
        cleaned = "processed"
    return cleaned[:60]


def _generate_pdf_plan(*, openai_client: OpenAI, model: str, prompt: str) -> Dict[str, str]:
    system_prompt = (
        "You are a PDF workflow planner.\n"
        "Analyze user intent, select the skill, and generate a runnable Python script.\n"
        "Output strict JSON only with schema:\n"
        "{\n"
        '  "analysis": "short intent analysis",\n'
        '  "selected_skill": "pdf",\n'
        '  "python_script": "python code string",\n'
        '  "assistant_markdown": "user-facing explanation"\n'
        "}\n"
        "Rules:\n"
        "- selected_skill must be 'pdf'.\n"
        "- python_script must use fitz (PyMuPDF) only.\n"
        "- python_script must read INPUT_PDF and OUTPUT_PDF from os.environ.\n"
        "- python_script must write a valid output PDF to OUTPUT_PDF.\n"
        "- python_script should perform a safe, visible modification that matches the user request.\n"
        "- python_script must print one JSON object to stdout: {\"summary\": \"...\"}.\n"
    )
    user_prompt = (
        f"User request:\n{prompt}\n\n"
        "Generate a robust script. If request is ambiguous, do a conservative edit: add a footer watermark "
        "with processing timestamp to each page."
    )
    completion = openai_client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.0,
    )
    raw = completion.choices[0].message.content or ""
    parsed = _extract_json_object(raw)
    script = _extract_script(str(parsed.get("python_script", "")))
    analysis = str(parsed.get("analysis", "")).strip()
    selected_skill = str(parsed.get("selected_skill", "pdf")).strip().lower() or "pdf"
    assistant_markdown = str(parsed.get("assistant_markdown", "")).strip()

    if not script:
        script = (
            "import json\n"
            "import os\n"
            "from datetime import datetime\n"
            "import fitz\n\n"
            "input_pdf = os.environ['INPUT_PDF']\n"
            "output_pdf = os.environ['OUTPUT_PDF']\n\n"
            "doc = fitz.open(input_pdf)\n"
            "stamp = datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')\n"
            "for page in doc:\n"
            "    page.insert_text((36, page.rect.height - 24), f'Processed by MatOpt ({stamp})', fontsize=9)\n"
            "doc.save(output_pdf)\n"
            "print(json.dumps({'summary': 'Added processing footer watermark to each page.'}))\n"
        )
    if not analysis:
        analysis = "Detected a PDF modification request and selected the PDF skill."
    if selected_skill != "pdf":
        selected_skill = "pdf"
    if not assistant_markdown:
        assistant_markdown = (
            "I analyzed your request, selected the PDF skill, executed a Python script, and generated a new PDF."
        )

    return {
        "analysis": analysis,
        "selected_skill": selected_skill,
        "python_script": script,
        "assistant_markdown": assistant_markdown,
    }


def _run_pdf_script(*, source_bytes: bytes, script: str, user_prompt: str) -> Dict[str, Any]:
    with TemporaryDirectory(prefix="matopt_pdf_workflow_") as tmp:
        tmp_dir = Path(tmp)
        input_pdf = tmp_dir / "input.pdf"
        output_pdf = tmp_dir / "output.pdf"
        script_path = tmp_dir / "workflow.py"

        input_pdf.write_bytes(source_bytes)
        script_path.write_text(script, encoding="utf-8")

        env = {
            "PATH": os.getenv("PATH", ""),
            "PYTHONUNBUFFERED": "1",
            "INPUT_PDF": str(input_pdf),
            "OUTPUT_PDF": str(output_pdf),
            "USER_PROMPT": user_prompt,
        }

        proc = subprocess.run(
            ["python3", str(script_path)],
            cwd=str(tmp_dir),
            env=env,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=45,
            check=False,
        )

        if proc.returncode != 0:
            detail = (proc.stderr or proc.stdout or "").strip()
            raise HTTPException(status_code=500, detail=f"PDF workflow script failed: {detail[:1400]}")
        if not output_pdf.exists() or output_pdf.stat().st_size == 0:
            raise HTTPException(status_code=500, detail="PDF workflow script did not generate output.")

        summary = ""
        parsed_stdout = _extract_json_object(proc.stdout or "")
        if parsed_stdout:
            summary = str(parsed_stdout.get("summary", "")).strip()

        return {
            "output_bytes": output_pdf.read_bytes(),
            "summary": summary,
        }


def _persist_generated_pdf(*, original_filename: str, content: bytes) -> WorkflowDownload:
    stem = Path(original_filename).stem or "processed"
    safe_stem = _safe_output_filename(stem)
    token = uuid.uuid4().hex[:10]
    final_name = f"{safe_stem}-modified-{token}.pdf"
    final_path = GENERATED_ROOT / final_name
    final_path.write_bytes(content)
    return WorkflowDownload(
        url=f"/generated/{final_name}",
        filename=final_name,
        size_bytes=len(content),
    )


def _run_pdf_workflow(*, prompt: str, model: str, source_filename: str, source_bytes: bytes) -> PdfWorkflowResponse:
    if fitz is None:
        raise HTTPException(
            status_code=500,
            detail="PyMuPDF is required for PDF workflows. Install 'pymupdf'.",
        )

    client = get_client()
    routing = resolve_skill_selection(
        messages=[{"role": "user", "content": prompt}],
        openai_client=None,
        model=None,
    )
    routed_skills = [str(s) for s in routing.get("selected_slugs", []) if str(s).strip()]
    if "pdf" not in routed_skills:
        routed_skills = ["pdf", *[s for s in routed_skills if s != "pdf"]]

    plan = _generate_pdf_plan(openai_client=client, model=model, prompt=prompt)
    run_result = _run_pdf_script(
        source_bytes=source_bytes,
        script=str(plan.get("python_script", "")),
        user_prompt=prompt,
    )
    download = _persist_generated_pdf(
        original_filename=source_filename,
        content=run_result["output_bytes"],
    )
    summary = str(run_result.get("summary", "")).strip()

    assistant_markdown = str(plan.get("assistant_markdown", "")).strip()
    route_explanation = str(routing.get("explanation", "")).strip()
    skill_list = ", ".join(f"`{s}`" for s in routed_skills[:4]) or "`pdf`"
    header = (
        "Skill routing:\n"
        f"- Selected skill(s): {skill_list}\n"
        f"- Prompt analysis: {str(plan.get('analysis', '')).strip()}\n"
    )
    if route_explanation:
        header += f"- Router rationale: {route_explanation}\n"
    assistant_markdown = f"{header}\n{assistant_markdown}".strip()
    if summary:
        assistant_markdown = f"{assistant_markdown}\n\nWhat I changed: {summary}"

    return PdfWorkflowResponse(
        assistant_markdown=assistant_markdown,
        analysis=str(plan.get("analysis", "")),
        selected_skills=routed_skills[:4],
        download=download,
        metadata={
            "workflow": "pdf",
            "selected_skills": routed_skills[:4],
            "router_strategy": str(routing.get("strategy", "")),
        },
    )


def _render_text_preview_image(text: str) -> str:
    width, height = 1200, 1600
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()

    margin = 56
    y = 56
    max_width_chars = 94

    for line in text.splitlines()[:70]:
        wrapped_lines = textwrap.wrap(line, width=max_width_chars) or [""]
        for wrapped in wrapped_lines:
            if y > height - margin - 18:
                break
            draw.text((margin, y), wrapped, fill="#222222", font=font)
            y += 18
        if y > height - margin - 18:
            break

    buf = io.BytesIO()
    image.save(buf, format="PNG")
    return _to_data_url("image/png", buf.getvalue())


def _pdf_to_markdown_and_preview(pdf_bytes: bytes) -> tuple[str, Optional[str], int]:
    if fitz is None:
        raise HTTPException(
            status_code=500,
            detail="PyMuPDF is required for PDF parsing. Install 'pymupdf'.",
        )

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    blocks: List[str] = []

    for idx, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            blocks.append(f"## Page {idx}\n\n{text}")

    preview: Optional[str] = None
    if doc.page_count > 0:
        first_page = doc.load_page(0)
        pix = first_page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
        preview = _to_data_url("image/png", pix.tobytes("png"))

    markdown = "\n\n".join(blocks).strip() or "_No extractable text found in PDF._"
    return markdown, preview, doc.page_count


def _docx_to_markdown_and_xml(data: bytes) -> tuple[str, str]:
    if Document is None:
        raise HTTPException(
            status_code=500,
            detail="python-docx is required for DOCX parsing. Install 'python-docx'.",
        )

    doc = Document(io.BytesIO(data))
    md_lines: List[str] = []

    root = ET.Element("document", attrib={"type": "docx"})

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        style_name = paragraph.style.name if paragraph.style else "Normal"
        xml_p = ET.SubElement(root, "paragraph", attrib={"style": style_name})
        xml_p.text = text

        if style_name.startswith("Heading"):
            level = 1
            digits = "".join(ch for ch in style_name if ch.isdigit())
            if digits:
                level = min(max(int(digits), 1), 6)
            md_lines.append(f"{'#' * level} {text}")
        else:
            md_lines.append(text)

    for table_idx, table in enumerate(doc.tables, start=1):
        xml_table = ET.SubElement(root, "table", attrib={"index": str(table_idx)})
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            xml_row = ET.SubElement(xml_table, "row")
            for cell_text in cells:
                xml_cell = ET.SubElement(xml_row, "cell")
                xml_cell.text = cell_text
            rows.append(cells)

        if rows:
            header = " | ".join(rows[0])
            separator = " | ".join(["---"] * len(rows[0]))
            md_lines.append(f"\n| {header} |\n| {separator} |")
            for row in rows[1:]:
                md_lines.append(f"| {' | '.join(row)} |")

    markdown = "\n\n".join(md_lines).strip() or "_No extractable text found in DOCX._"
    xml_string = ET.tostring(root, encoding="unicode")
    return markdown, xml_string


def _pptx_to_markdown_and_xml(data: bytes) -> tuple[str, str]:
    if Presentation is None:
        raise HTTPException(
            status_code=500,
            detail="python-pptx is required for PPTX parsing. Install 'python-pptx'.",
        )

    presentation = Presentation(io.BytesIO(data))
    sections: List[str] = []

    root = ET.Element("presentation", attrib={"type": "pptx"})

    for index, slide in enumerate(presentation.slides, start=1):
        xml_slide = ET.SubElement(root, "slide", attrib={"index": str(index)})

        title_text = ""
        if slide.shapes.title and slide.shapes.title.text:
            title_text = slide.shapes.title.text.strip()

        if title_text:
            sections.append(f"## Slide {index}: {title_text}")
            ET.SubElement(xml_slide, "title").text = title_text
        else:
            sections.append(f"## Slide {index}")

        bullet_lines: List[str] = []
        for shape in slide.shapes:
            shape_text = getattr(shape, "text", "")
            if not shape_text:
                continue
            text = shape_text.strip()
            if not text or text == title_text:
                continue
            bullet_lines.append(f"- {text.replace(chr(10), ' ')}")
            ET.SubElement(xml_slide, "text").text = text

        if bullet_lines:
            sections.append("\n".join(bullet_lines))
        else:
            sections.append("_No text found on this slide._")

    markdown = "\n\n".join(sections).strip() or "_No extractable text found in PPTX._"
    xml_string = ET.tostring(root, encoding="unicode")
    return markdown, xml_string


def _ipynb_to_markdown_and_xml(data: bytes) -> tuple[str, str]:
    text = _safe_decode(data)
    notebook = json.loads(text)
    cells = notebook.get("cells", [])

    md_blocks: List[str] = []
    root = ET.Element("notebook", attrib={"type": "ipynb"})

    for index, cell in enumerate(cells, start=1):
        cell_type = cell.get("cell_type", "unknown")
        source = "".join(cell.get("source", []))

        xml_cell = ET.SubElement(
            root,
            "cell",
            attrib={"index": str(index), "type": str(cell_type)},
        )
        ET.SubElement(xml_cell, "source").text = source

        if cell_type == "markdown":
            md_blocks.append(source.strip())
        elif cell_type == "code":
            md_blocks.append(f"```python\n{source.rstrip()}\n```")

            outputs = cell.get("outputs", [])
            rendered_outputs: List[str] = []
            for output in outputs:
                text_output = output.get("text")
                if isinstance(text_output, list):
                    rendered_outputs.append("".join(text_output).strip())
                elif isinstance(text_output, str):
                    rendered_outputs.append(text_output.strip())

            if rendered_outputs:
                joined_output = "\n".join(x for x in rendered_outputs if x)
                if joined_output:
                    md_blocks.append(f"Output:\n```text\n{joined_output}\n```")
                    ET.SubElement(xml_cell, "output").text = joined_output

    markdown = "\n\n".join(block for block in md_blocks if block).strip()
    if not markdown:
        markdown = "_Notebook parsed, but no readable cell content was found._"

    xml_string = ET.tostring(root, encoding="unicode")
    return markdown, xml_string


def _try_convert_office_to_pdf(data: bytes, ext: str) -> Optional[bytes]:
    soffice_path = shutil.which("soffice")
    if not soffice_path:
        return None

    with TemporaryDirectory() as temp_dir:
        in_path = Path(temp_dir) / f"input.{ext}"
        in_path.write_bytes(data)

        cmd = [
            soffice_path,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            temp_dir,
            str(in_path),
        ]

        try:
            subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=90,
            )
        except Exception:
            return None

        pdf_candidates = list(Path(temp_dir).glob("*.pdf"))
        if not pdf_candidates:
            return None
        return pdf_candidates[0].read_bytes()


def _text_to_markdown(data: bytes, language: str) -> str:
    source = _safe_decode(data)
    return f"```{language}\n{source.rstrip()}\n```"


def _run_chat(payload: ChatRequest) -> ChatResponse:
    if not payload.messages:
        raise HTTPException(status_code=400, detail="messages cannot be empty")

    model = payload.model or os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    api_key = _get_openai_api_key()
    client = get_client()
    message_dicts = [m.model_dump() for m in payload.messages]
    skill_selection = resolve_skill_selection(
        messages=message_dicts,
        openai_client=client,
        model=model,
        explicit_skill_slugs=payload.skill_slugs,
    )
    selected_skill_slugs = [
        str(slug) for slug in skill_selection.get("selected_slugs", []) if str(slug).strip()
    ]

    if payload.agentic is not False:
        try:
            agentic_result = run_agentic_chat(
                openai_client=client,
                model=model,
                api_key=api_key,
                messages=message_dicts,
                skill_slugs=selected_skill_slugs,
            )
            return ChatResponse(
                id=f"chatcmpl-{int(time.time() * 1000)}",
                content=str(agentic_result.get("assistant_markdown", "")),
                artifacts=[
                    VisualArtifact.model_validate(item)
                    for item in agentic_result.get("artifacts", [])
                    if isinstance(item, dict)
                ],
                tool_calls=[
                    item
                    for item in agentic_result.get("tool_calls", [])
                    if isinstance(item, dict)
                ],
                metadata={
                    "agentic": True,
                    "agent_mode": agentic_result.get("agent_mode", "langgraph"),
                    "selected_skills": selected_skill_slugs,
                },
            )
        except Exception as exc:
            raise HTTPException(
                status_code=500,
                detail=f"Agentic pipeline failed: {exc}",
            ) from exc

    try:
        completion = client.chat.completions.create(
            model=model,
            messages=message_dicts,
            temperature=0.7,
        )
    except APIStatusError as exc:
        status = exc.status_code or 502
        raise HTTPException(
            status_code=status,
            detail=f"OpenAI API error ({status}): {exc}",
        ) from exc
    except APITimeoutError as exc:
        raise HTTPException(status_code=504, detail=f"OpenAI timeout: {exc}") from exc
    except APIConnectionError as exc:
        raise HTTPException(status_code=502, detail=f"OpenAI connection error: {exc}") from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"OpenAI request failed: {exc}") from exc

    content = completion.choices[0].message.content or ""
    usage = None
    if completion.usage:
        try:
            usage = completion.usage.model_dump()  # type: ignore[attr-defined]
        except Exception:
            usage = {"raw_usage": str(completion.usage)}

    return ChatResponse(
        id=f"chatcmpl-{int(time.time() * 1000)}",
        content=content,
        metadata={
            "agentic": False,
            "selected_skills": selected_skill_slugs,
        },
        usage=usage,
    )


def _stream_chat_sse(payload: ChatRequest):
    if not payload.messages:
        raise HTTPException(status_code=400, detail="messages cannot be empty")

    model = payload.model or os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    api_key = _get_openai_api_key()
    client = get_client()
    message_dicts = [m.model_dump() for m in payload.messages]
    skill_selection = resolve_skill_selection(
        messages=message_dicts,
        openai_client=client,
        model=model,
        explicit_skill_slugs=payload.skill_slugs,
    )
    selected_skill_slugs = [
        str(slug) for slug in skill_selection.get("selected_slugs", []) if str(slug).strip()
    ]

    if payload.agentic is not False:
        try:
            agentic = run_agentic_chat(
                openai_client=client,
                model=model,
                api_key=api_key,
                messages=message_dicts,
                skill_slugs=selected_skill_slugs,
            )
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"Agentic pipeline failed: {exc}") from exc

        content = str(agentic.get("assistant_markdown", ""))
        artifacts = [item for item in agentic.get("artifacts", []) if isinstance(item, dict)]
        tool_calls = [item for item in agentic.get("tool_calls", []) if isinstance(item, dict)]

        def agentic_event_gen():
            response_id = f"chatcmpl-{int(time.time() * 1000)}"
            try:
                step = 1
                for i in range(0, len(content), step):
                    yield f"data: {json.dumps({'delta': content[i:i + step]})}\n\n"

                done_payload = {
                    "done": True,
                    "id": response_id,
                    "artifacts": artifacts,
                    "tool_calls": tool_calls,
                    "metadata": {
                        "agentic": True,
                        "agent_mode": agentic.get("agent_mode", "langgraph"),
                        "selected_skills": selected_skill_slugs,
                    },
                }
                yield f"data: {json.dumps(done_payload)}\n\n"
                yield "data: [DONE]\n\n"
            except Exception as exc:
                err = str(exc)
                yield f"data: {json.dumps({'error': err})}\n\n"
                yield "data: [DONE]\n\n"

        return StreamingResponse(
            agentic_event_gen(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",
            },
        )

    try:
        stream = client.chat.completions.create(
            model=model,
            messages=message_dicts,
            temperature=0.7,
            stream=True,
        )
    except APIStatusError as exc:
        status = exc.status_code or 502
        raise HTTPException(
            status_code=status,
            detail=f"OpenAI API error ({status}): {exc}",
        ) from exc
    except APITimeoutError as exc:
        raise HTTPException(status_code=504, detail=f"OpenAI timeout: {exc}") from exc
    except APIConnectionError as exc:
        raise HTTPException(status_code=502, detail=f"OpenAI connection error: {exc}") from exc
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"OpenAI request failed: {exc}") from exc

    def event_gen():
        response_id = f"chatcmpl-{int(time.time() * 1000)}"
        try:
            for chunk in stream:
                delta = ""
                try:
                    if chunk.choices and chunk.choices[0].delta:
                        delta = chunk.choices[0].delta.content or ""
                except Exception:
                    delta = ""

                if delta:
                    yield f"data: {json.dumps({'delta': delta})}\n\n"

            yield f"data: {json.dumps({'done': True, 'id': response_id, 'metadata': {'agentic': False, 'selected_skills': selected_skill_slugs}})}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as exc:
            err = str(exc)
            yield f"data: {json.dumps({'error': err})}\n\n"
            yield "data: [DONE]\n\n"

    return StreamingResponse(
        event_gen(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.get("/")
def root() -> dict:
    return {"service": "matopt-backend", "status": "ok"}


@app.get("/health")
def health() -> dict:
    return {"ok": True}


@app.post("/chat", response_model=ChatResponse)
def chat(payload: ChatRequest) -> ChatResponse:
    return _run_chat(payload)


@app.post("/api/chat", response_model=ChatResponse)
def chat_api(payload: ChatRequest) -> ChatResponse:
    return _run_chat(payload)


@app.post("/chat/stream")
def chat_stream(payload: ChatRequest):
    return _stream_chat_sse(payload)


@app.post("/api/chat/stream")
def chat_stream_api(payload: ChatRequest):
    return _stream_chat_sse(payload)


@app.post("/workflows/pdf", response_model=PdfWorkflowResponse)
async def workflows_pdf(
    prompt: str = Form(...),
    file: UploadFile = File(...),
    model: Optional[str] = Form(None),
) -> PdfWorkflowResponse:
    if not file.filename:
        raise HTTPException(status_code=400, detail="Filename is required")
    ext = _get_extension(file.filename)
    if ext != "pdf":
        raise HTTPException(status_code=400, detail="This workflow currently supports PDF files only.")

    payload = (prompt or "").strip()
    if not payload:
        raise HTTPException(status_code=400, detail="Prompt is required.")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    resolved_model = model or os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    return _run_pdf_workflow(
        prompt=payload,
        model=resolved_model,
        source_filename=file.filename,
        source_bytes=data,
    )


@app.post("/api/workflows/pdf", response_model=PdfWorkflowResponse)
async def workflows_pdf_api(
    prompt: str = Form(...),
    file: UploadFile = File(...),
    model: Optional[str] = Form(None),
) -> PdfWorkflowResponse:
    return await workflows_pdf(prompt=prompt, file=file, model=model)


def _skill_response(data: Dict[str, Any]) -> SkillResponse:
    return SkillResponse(
        slug=str(data["slug"]),
        title=str(data["title"]),
        content=str(data["content"]),
        path=str(data["path"]),
        updated_at=str(data["updated_at"]),
        size_bytes=int(data["size_bytes"]),
    )


@app.get("/skills")
def skills_list() -> Dict[str, Any]:
    items = list_skills()
    return {"skills": items, "count": len(items), "generated_at": datetime.now(timezone.utc).isoformat()}


@app.get("/api/skills")
def skills_list_api() -> Dict[str, Any]:
    return skills_list()


@app.get("/skills/{slug}", response_model=SkillResponse)
def skills_get(slug: str) -> SkillResponse:
    try:
        return _skill_response(read_skill(slugify(slug)))
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/api/skills/{slug}", response_model=SkillResponse)
def skills_get_api(slug: str) -> SkillResponse:
    return skills_get(slug)


@app.post("/skills", response_model=SkillResponse)
def skills_upsert(payload: SkillUpsertRequest) -> SkillResponse:
    data = upsert_skill(
        slug=payload.slug,
        content=payload.content,
        title=payload.title,
    )
    return _skill_response(data)


@app.post("/api/skills", response_model=SkillResponse)
def skills_upsert_api(payload: SkillUpsertRequest) -> SkillResponse:
    return skills_upsert(payload)


@app.delete("/skills/{slug}")
def skills_delete(slug: str) -> Dict[str, Any]:
    deleted = delete_skill(slugify(slug))
    if not deleted:
        raise HTTPException(status_code=404, detail=f"Skill '{slug}' does not exist")
    return {"deleted": True, "slug": slugify(slug)}


@app.delete("/api/skills/{slug}")
def skills_delete_api(slug: str) -> Dict[str, Any]:
    return skills_delete(slug)


def _parse_document_impl(filename: str, data: bytes) -> DocumentParseResponse:
    ext = _get_extension(filename)
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=(
                "Unsupported file type. Allowed: "
                + ", ".join(sorted(SUPPORTED_EXTENSIONS))
            ),
        )

    metadata: Dict[str, Any] = {
        "filename": filename,
        "extension": ext,
        "size_bytes": len(data),
        "preview_source": None,
        "structured_xml": None,
    }

    markdown = ""
    preview_image_data_url: Optional[str] = None

    if ext == "pdf":
        markdown, preview_image_data_url, page_count = _pdf_to_markdown_and_preview(data)
        metadata["page_count"] = page_count
        metadata["preview_source"] = "pdf_page_1"

    elif ext in IMAGE_EXTENSIONS:
        markdown = f"# Image: {filename}\n\nUploaded image file."
        preview_image_data_url = _to_data_url(_mime_for_ext(ext), data)
        metadata["preview_source"] = "original_image"

    elif ext in TEXT_EXTENSIONS:
        markdown = _text_to_markdown(data, "r" if ext in {"r", "rmd"} else ext)
        preview_image_data_url = _render_text_preview_image(_safe_decode(data))
        metadata["preview_source"] = "rendered_text_preview"

    elif ext == "ipynb":
        markdown, xml = _ipynb_to_markdown_and_xml(data)
        preview_image_data_url = _render_text_preview_image(markdown)
        metadata["structured_xml"] = xml
        metadata["preview_source"] = "rendered_text_preview"

    elif ext == "docx":
        markdown, xml = _docx_to_markdown_and_xml(data)
        metadata["structured_xml"] = xml

        office_pdf = _try_convert_office_to_pdf(data, ext)
        if office_pdf:
            _, preview_image_data_url, _ = _pdf_to_markdown_and_preview(office_pdf)
            metadata["preview_source"] = "libreoffice_pdf_page_1"
        else:
            preview_image_data_url = _render_text_preview_image(markdown)
            metadata["preview_source"] = "rendered_text_preview"
            metadata["preview_note"] = (
                "Install LibreOffice to generate true first-page previews for DOCX/PPTX/PPT/DOC."
            )

    elif ext == "pptx":
        markdown, xml = _pptx_to_markdown_and_xml(data)
        metadata["structured_xml"] = xml

        office_pdf = _try_convert_office_to_pdf(data, ext)
        if office_pdf:
            _, preview_image_data_url, _ = _pdf_to_markdown_and_preview(office_pdf)
            metadata["preview_source"] = "libreoffice_pdf_page_1"
        else:
            preview_image_data_url = _render_text_preview_image(markdown)
            metadata["preview_source"] = "rendered_text_preview"
            metadata["preview_note"] = (
                "Install LibreOffice to generate true first-page previews for DOCX/PPTX/PPT/DOC."
            )

    elif ext in {"doc", "ppt"}:
        office_pdf = _try_convert_office_to_pdf(data, ext)
        if office_pdf:
            markdown, preview_image_data_url, page_count = _pdf_to_markdown_and_preview(office_pdf)
            metadata["page_count"] = page_count
            metadata["preview_source"] = "libreoffice_pdf_page_1"
            metadata["conversion"] = "binary_office_to_pdf"
        else:
            markdown = (
                f"# {filename}\n\n"
                "Binary Office format detected. Text extraction needs LibreOffice in the backend runtime."
            )
            preview_image_data_url = _render_text_preview_image(markdown)
            metadata["preview_source"] = "rendered_text_preview"
            metadata["conversion"] = "fallback_without_libreoffice"
            metadata["preview_note"] = (
                "Install LibreOffice to parse DOC/PPT and generate true first-page previews."
            )

    else:
        raise HTTPException(status_code=400, detail="Unsupported file type.")

    return DocumentParseResponse(
        markdown=markdown,
        metadata=metadata,
        preview_image_data_url=preview_image_data_url,
    )


@app.post("/documents/parse", response_model=DocumentParseResponse)
async def parse_document(file: UploadFile = File(...)) -> DocumentParseResponse:
    if not file.filename:
        raise HTTPException(status_code=400, detail="Filename is required")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    return _parse_document_impl(file.filename, data)


@app.post("/api/documents/parse", response_model=DocumentParseResponse)
async def parse_document_api(file: UploadFile = File(...)) -> DocumentParseResponse:
    return await parse_document(file)


@app.get("/api/health")
def api_health() -> dict:
    return {"ok": True}
